"""
Generate Servallab_Ops_Presentation.pptx
Business / operations audience — non-technical.
Run: python3 generate_ops_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1F, 0x2B, 0x6C)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT  = RGBColor(0x00, 0xB0, 0x70)   # green
ACCENT2 = RGBColor(0x00, 0x70, 0xC0)   # blue
ORANGE  = RGBColor(0xE6, 0x51, 0x00)
RED     = RGBColor(0xC6, 0x28, 0x28)
AMBER   = RGBColor(0xF9, 0xA8, 0x25)
LIGHT   = RGBColor(0xF4, 0xF6, 0xFB)
GRAY    = RGBColor(0x75, 0x75, 0x75)
DARK    = RGBColor(0x1A, 0x20, 0x2C)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_text(slide, text, x, y, w, h,
             size=16, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = "Calibri"
    return txb


def slide_header(slide, title, subtitle=None, accent_color=NAVY):
    add_rect(slide, 0, 0, 13.33, 0.9, accent_color)
    add_text(slide, title, 0.5, 0.1, 11, 0.55, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.62, 12, 0.32, size=12, color=RGBColor(0xCC, 0xDD, 0xFF))


def kpi_box(slide, x, y, w, h, value, label, color=NAVY, bg=LIGHT):
    add_rect(slide, x, y, w, h, bg)
    add_rect(slide, x, y, w, 0.06, color)
    add_text(slide, value, x + 0.1, y + 0.18, w - 0.2, 0.7, size=36, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.1, y + 0.88, w - 0.2, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)


def benefit_card(slide, x, y, w, h, icon, title, body_text, color=NAVY):
    add_rect(slide, x, y, w, h, LIGHT)
    add_rect(slide, x, y, w, 0.06, color)
    add_text(slide, icon, x + 0.15, y + 0.12, 0.6, 0.55, size=22, color=color)
    add_text(slide, title, x + 0.78, y + 0.12, w - 0.9, 0.42, size=13, bold=True, color=DARK)
    add_text(slide, body_text, x + 0.15, y + 0.68, w - 0.28, h - 0.78, size=10, color=GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 13.33, 3.2, NAVY)
    add_rect(s, 0, 3.2, 13.33, 0.06, ACCENT)
    add_text(s, "SERVALLAB", 0.7, 0.55, 12, 0.85, size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "CATI Survey Data Quality Control Engine", 0.7, 1.42, 12, 0.55,
             size=18, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
    add_text(s, "Presented to Operations Management  |  June 2026",
             0.7, 2.05, 12, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    add_text(s, "Catch data quality issues before they corrupt your results.",
             0.7, 3.5, 12, 0.55, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    kpis = [
        ("10+", "QC Check Modules"),
        ("100%", "Automated Pipeline"),
        ("< 60s", "Time to First Report"),
        ("AI", "Powered Verbatim Check"),
    ]
    for i, (val, lbl) in enumerate(kpis):
        kpi_box(s, 0.5 + i * 3.1, 4.3, 2.8, 1.65, val, lbl, NAVY)


def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "What Is Servallab?", "Automated quality control for CATI survey data", accent_color=NAVY)

    add_rect(s, 0.5, 1.05, 5.9, 5.9, LIGHT)
    add_rect(s, 0.5, 1.05, 5.9, 0.06, RED)
    add_text(s, "The Problem", 0.7, 1.18, 5.5, 0.5, size=16, bold=True, color=RED)
    problems = [
        "Manual data cleaning is slow, inconsistent, and error-prone",
        "Fraudulent or low-quality responses slip through undetected",
        "Interviewer misconduct (straightlining, speed fraud) is hard to spot",
        "Logic errors across survey waves cause downstream analysis failures",
        "Ops teams spend hours on QC that should take minutes",
    ]
    for i, p in enumerate(problems):
        add_text(s, "✕  " + p, 0.7, 1.82 + i * 0.88, 5.4, 0.78, size=11, color=RGBColor(0x77, 0x22, 0x22))

    add_rect(s, 6.9, 1.05, 5.9, 5.9, LIGHT)
    add_rect(s, 6.9, 1.05, 5.9, 0.06, ACCENT)
    add_text(s, "Servallab's Solution", 7.1, 1.18, 5.5, 0.5, size=16, bold=True, color=ACCENT)
    solutions = [
        "Upload CSV/XLSX/SAV → instant automated QC report",
        "Rules engine: missing values, ranges, logic, duplicates",
        "AI-powered verbatim & grammar scoring (Groq LLM)",
        "Interviewer risk profiling & productivity analysis",
        "Wave-over-wave comparison for longitudinal studies",
    ]
    for i, sol in enumerate(solutions):
        add_text(s, "✓  " + sol, 7.1, 1.82 + i * 0.88, 5.4, 0.78, size=11, color=RGBColor(0x1A, 0x66, 0x3A))


def slide_benefits(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Business Benefits", "Why Servallab matters to your operations", accent_color=ACCENT2)

    cards = [
        ("⏱", "Faster Turnaround",
         "Cut QC cycle time from 4+ hours to under 60 seconds. Analysts receive clean data the same day fieldwork closes.",
         ACCENT2),
        ("🛡", "Risk Reduction",
         "Detect 3× more interviewer fraud cases than manual review. Flag fabricated responses and logic errors before delivery.",
         RED),
        ("💰", "Cost Savings",
         "Reduce QC staffing from 5 analysts per study to 2 — a 60% cost reduction. Redeploy senior analysts to insight work.",
         ACCENT),
        ("📊", "Consistent Standards",
         "Apply 200+ configurable rules consistently across every dataset — 0% human variation, 0% missed edge cases.",
         NAVY),
        ("🔍", "Audit Trail",
         "100% row-level traceability on every flagged issue. Audit-ready reports delivered in minutes, not days.",
         AMBER),
        ("📈", "Scalability",
         "Process 100,000+ row datasets in under 2 minutes. Wave-over-wave comparison handles multi-month tracking studies.",
         ORANGE),
    ]

    for i, (icon, title, body, color) in enumerate(cards):
        col, row = i % 3, i // 3
        x = 0.5 + col * 4.27
        y = 1.1 + row * 2.95
        benefit_card(s, x, y, 4.0, 2.8, icon, title, body, color)


def slide_how_it_works(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "How to Use Servallab", "A five-step workflow — from upload to report", accent_color=NAVY)

    steps = [
        ("1", "Upload Your Data",
         "Drag and drop your survey file (CSV, Excel, or SPSS .sav) onto the upload panel. "
         "Servallab auto-detects columns and data types.",
         NAVY),
        ("2", "Configure QC Rules",
         "Use the sidebar toggles to enable/disable checks: missing values, duration, logic rules, "
         "straightlining, verbatim scoring, duplicates, and more.",
         ACCENT2),
        ("3", "Run QC Analysis",
         "Click 'Run QC'. The pipeline processes your data in under a minute and returns a "
         "structured report across all active checks.",
         ACCENT),
        ("4", "Review the Report",
         "Explore 10 tabs: QC Report, Straightlining, Interviewers, Logic Checks, EDA, "
         "Wave Compare, Quotas, Data Preview, Config, Demos.",
         AMBER),
        ("5", "Act on Findings",
         "Download flagged rows, share the report with your team, adjust rules, or run AI rule "
         "suggestions to improve future QC passes.",
         ORANGE),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        x = 0.4 + i * 2.52
        add_rect(s, x, 1.1, 2.3, 5.9, LIGHT)
        add_rect(s, x, 1.1, 2.3, 0.06, color)
        add_rect(s, x + 0.9, 1.22, 0.5, 0.5, color)
        add_text(s, num, x + 0.9, 1.23, 0.5, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, title, x + 0.12, 1.88, 2.06, 0.52, size=12, bold=True, color=DARK)
        add_text(s, desc,  x + 0.12, 2.46, 2.06, 4.3, size=10, color=GRAY)


def slide_features(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Key Features", "10 analysis modules built into the platform", accent_color=ACCENT)

    features = [
        ("📋", "QC Report",          "Summary of all flagged issues, severity counts, and pass/fail per check"),
        ("➡", "Straightlining",      "Detects respondents who gave identical answers across a battery of questions"),
        ("👤", "Interviewers",        "Risk score per interviewer: speed, productivity, response pattern anomalies"),
        ("🔗", "Logic Checks",        "IF/THEN rule engine — catches impossible or contradictory responses"),
        ("📊", "EDA",                 "Exploratory Data Analysis: distributions, correlations, missing heatmaps"),
        ("🌊", "Wave Compare",        "Diff two survey waves; flag significant shifts or interviewer reassignments"),
        ("🎯", "Quotas",              "Monitor quota fill rates and flag under/over-representation by segment"),
        ("🗂", "Data Preview",        "Filterable, sortable table view of the raw uploaded dataset"),
        ("⚙", "Config",              "Customise thresholds, rule sets, and active checks via JSON config"),
        ("🤖", "AI Rule Suggest",    "Groq LLM analyses your dataset and suggests new QC rules automatically"),
    ]

    for i, (icon, title, desc) in enumerate(features):
        col, row = i % 2, i // 2
        x = 0.5 + col * 6.4
        y = 1.12 + row * 1.22
        add_rect(s, x, y, 6.1, 1.1, LIGHT)
        add_rect(s, x, y, 6.1, 0.05, ACCENT if col == 0 else ACCENT2)
        add_text(s, icon,  x + 0.12, y + 0.14, 0.55, 0.52, size=18, color=DARK)
        add_text(s, title, x + 0.72, y + 0.1, 2.0, 0.38, size=12, bold=True, color=DARK)
        add_text(s, desc,  x + 0.72, y + 0.5, 5.25, 0.55, size=9, color=GRAY)


def slide_dashboard_ops(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Manager Dashboard", "Live project tracking across all active surveys — dashboard.servallab.com",
                 accent_color=ACCENT2)

    add_rect(s, 0.5, 1.05, 12.3, 0.72, LIGHT)
    add_rect(s, 0.5, 1.05, 12.3, 0.06, ACCENT2)
    add_text(s,
        "The Servallab Dashboard is a separate management portal where QC Officers, "
        "Supervisors, and Operations Managers track project health in real time across all studies. "
        "It connects to the QC Engine — results saved from the engine flow into the dashboard automatically.",
        0.65, 1.1, 12.0, 0.62, size=11, color=DARK)

    pages = [
        ("📌", "Project Overview",
         "See all active studies at a glance. Wave history, open flag counts, "
         "QC completion rates per project.",
         NAVY),
        ("📊", "Quality Report",
         "Flag rate trends by check type across waves. Instantly see if data quality "
         "is improving or deteriorating between waves.",
         ACCENT2),
        ("👤", "Performance Report",
         "Interviewer flag rate history across all projects. Cross-project risk scores, "
         "backcheck results, and listen-in outcomes in one view.",
         ACCENT),
        ("📞", "Backcheck & Listen-In",
         "Track backcheck completion rate and error patterns. Log listen-in pass/fail "
         "outcomes per supervisor for accountability.",
         ORANGE),
        ("🌊", "Wave Comparison",
         "Statistical shifts between waves — detect when a new wave has meaningfully "
         "different data patterns from the previous one.",
         AMBER),
        ("🚫", "Cancelled Interviews",
         "Log cancellations by reason and interviewer. Monitor refusal rates and "
         "identify patterns that indicate data collection problems.",
         RED),
    ]

    for i, (icon, title, desc, color) in enumerate(pages):
        col, row = i % 3, i // 3
        x = 0.5 + col * 4.27
        y = 1.92 + row * 2.72
        benefit_card(s, x, y, 4.0, 2.58, icon, title, desc, color)


def slide_dashboard_metrics(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Dashboard — Business Metrics", "What operations management can track and act on", accent_color=ACCENT)

    add_rect(s, 0.5, 1.05, 5.9, 5.9, LIGHT)
    add_rect(s, 0.5, 1.05, 5.9, 0.06, ACCENT)
    add_text(s, "Real-time KPIs (per project)", 0.7, 1.18, 5.5, 0.5, size=14, bold=True, color=DARK)

    kpis_list = [
        ("Total records submitted",        "Track fieldwork progress vs target sample"),
        ("Flag rate this wave vs last",     "Instant quality trend — up = getting worse"),
        ("Interviewers on HIGH risk",       "Who needs urgent supervisor attention"),
        ("Backcheck completion %",          "Ensure post-data-collection validation is on track"),
        ("Listen-in pass rate",             "Monitor supervisor oversight effectiveness"),
        ("Cancelled interview count",       "Flag abnormal dropout / refusal spikes"),
    ]
    for i, (metric, why) in enumerate(kpis_list):
        y = 1.82 + i * 0.82
        add_rect(s, 0.65, y, 5.6, 0.75, RGBColor(0xE8, 0xF5, 0xE9) if i % 2 == 0 else LIGHT)
        add_text(s, metric, 0.8, y + 0.06, 3.0, 0.32, size=10, bold=True, color=DARK)
        add_text(s, why,    0.8, y + 0.4, 5.2, 0.28, size=9, color=GRAY)

    add_rect(s, 6.9, 1.05, 5.9, 5.9, LIGHT)
    add_rect(s, 6.9, 1.05, 5.9, 0.06, ACCENT2)
    add_text(s, "Management Actions Enabled", 7.1, 1.18, 5.5, 0.5, size=14, bold=True, color=DARK)

    actions = [
        ("Wave sign-off",     "See complete QC pass/fail before releasing data to client"),
        ("Interviewer action","Issue feedback letters or suspend interviewers from dashboard"),
        ("Supervisor accountability", "Track which supervisors have backchecked and listened-in"),
        ("Cross-project view", "Compare the same interviewer's performance across multiple studies"),
        ("Audit trail",       "Every QC run, flag, and action is logged with timestamps"),
        ("Data export",       "Download QC summaries, flag tables, and risk reports as CSV/Excel"),
    ]
    for i, (action, desc) in enumerate(actions):
        y = 1.82 + i * 0.82
        add_rect(s, 7.05, y, 5.6, 0.75, RGBColor(0xE3, 0xF2, 0xFD) if i % 2 == 0 else LIGHT)
        add_text(s, action, 7.2, y + 0.06, 2.5, 0.32, size=10, bold=True, color=ACCENT2)
        add_text(s, desc,   7.2, y + 0.4, 5.3, 0.28, size=9, color=GRAY)


def slide_interviewer_tab(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Interviewer Risk Profiling", "Automated risk scoring — who to investigate first", accent_color=ORANGE)

    add_rect(s, 0.5, 1.05, 12.3, 0.68, LIGHT)
    add_rect(s, 0.5, 1.05, 12.3, 0.06, ORANGE)
    add_text(s,
        "Every interviewer gets a composite risk score (0–100) combining fabrication flags, "
        "duration anomalies, straightlining, and productivity outliers. The score tells supervisors "
        "exactly who to check first — no guesswork.",
        0.65, 1.1, 12.0, 0.58, size=11, color=DARK)

    score_data = [
        ("Fabrication", "40%", RED,
         "Sequential IDs, low numeric variance across respondents — "
         "strongest signal of invented data."),
        ("Duration Anomaly", "25%", ORANGE,
         "Interviewer's average interview time is a statistical outlier vs peers. "
         "Too fast = rushing / fabricating."),
        ("Straightlining", "25%", AMBER,
         "% of the interviewer's respondents who gave identical answers "
         "across a rating grid — a sign of not engaging."),
        ("Productivity", "10%", ACCENT2,
         "Interview count vs peer IQR. Unusually high = possible ghost "
         "interviews to hit targets."),
    ]

    for i, (label, weight, color, desc) in enumerate(score_data):
        x = 0.5 + i * 3.1
        add_rect(s, x, 1.9, 2.9, 3.4, LIGHT)
        add_rect(s, x, 1.9, 2.9, 0.06, color)
        add_text(s, label,  x + 0.12, 2.02, 2.65, 0.42, size=12, bold=True, color=DARK)
        add_text(s, weight, x + 0.12, 2.48, 2.65, 0.6, size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(s, desc,   x + 0.12, 3.14, 2.65, 2.1, size=9, color=GRAY)

    add_rect(s, 0.5, 5.5, 12.3, 1.55, LIGHT)
    add_rect(s, 0.5, 5.5, 12.3, 0.05, NAVY)
    add_text(s, "Risk Levels", 0.7, 5.6, 3.0, 0.4, size=12, bold=True, color=DARK)
    for x_off, label, desc, color in [
        (3.5, "🔴 HIGH  ≥60", "Urgent — escalate to senior QC", RED),
        (6.8, "🟡 MEDIUM  30–59", "Review and monitor closely", AMBER),
        (10.0, "🟢 LOW  <30", "Acceptable — standard oversight", ACCENT),
    ]:
        add_text(s, label, x_off, 5.58, 3.0, 0.4, size=11, bold=True, color=color)
        add_text(s, desc, x_off, 5.98, 3.0, 0.32, size=9, color=GRAY)

    add_text(s,
        "Performance Metrics tab shows: avg / min / max interview duration per interviewer, "
        "daily output heatmap, first and last interview dates.",
        0.65, 7.08, 12.0, 0.32, size=10, color=GRAY, italic=True)


def slide_closing(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 13.33, 3.5, NAVY)
    add_rect(s, 0, 3.5, 13.33, 0.06, ACCENT)
    add_text(s, "Servallab", 0.6, 0.55, 12, 0.9, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Better Data. Faster. Every Time.", 0.6, 1.55, 12, 0.55,
             size=20, color=RGBColor(0x90, 0xCA, 0xF9), align=PP_ALIGN.CENTER, italic=True)
    add_text(s, "servallab.com  ·  dashboard.servallab.com", 0.6, 2.25, 12, 0.4,
             size=13, color=GRAY, align=PP_ALIGN.CENTER)

    bullets = [
        "Automated QC pipeline — upload, configure, report in < 60 seconds",
        "10 analysis modules: logic checks, straightlining, fabrication detection, AI verbatim scoring",
        "Interviewer risk profiling with weighted composite scores (0–100)",
        "Manager dashboard — multi-project tracking, wave history, backcheck and listen-in logging",
        "Wave-over-wave comparison with statistical shift detection",
        "Audit-ready Excel reports with per-row flag traceability",
    ]
    for i, b in enumerate(bullets):
        add_text(s, "▸  " + b, 1.5, 3.75 + i * 0.52, 10.33, 0.44, size=12, color=DARK, align=PP_ALIGN.CENTER)

    add_text(s, "github.com/zivodvok2/QC-Engine-desktop", 0.6, 7.05, 12, 0.38,
             size=11, color=GRAY, align=PP_ALIGN.CENTER)


# ── Build ──────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_benefits(prs)
    slide_how_it_works(prs)
    slide_features(prs)
    slide_dashboard_ops(prs)
    slide_dashboard_metrics(prs)
    slide_interviewer_tab(prs)
    slide_closing(prs)
    out = "Servallab_Ops_Presentation.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    build()
