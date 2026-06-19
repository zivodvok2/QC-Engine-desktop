"""
Generate Servallab QC Engine — Technical Presentation
Run: python3 generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0E, 0x10, 0x14)
SURFACE  = RGBColor(0x16, 0x1A, 0x22)
ACCENT   = RGBColor(0x4A, 0xF0, 0xA0)
ACCENT2  = RGBColor(0x4A, 0x9E, 0xF0)
CRITICAL = RGBColor(0xF0, 0x5A, 0x5A)
WARNING  = RGBColor(0xF0, 0xC0, 0x4A)
INFO_COL = RGBColor(0x4A, 0xBF, 0xF0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x88, 0x96, 0xA8)
DARK     = RGBColor(0x1A, 0x20, 0x2C)
DARKER   = RGBColor(0x0A, 0x0C, 0x10)
DARKISH  = RGBColor(0x1E, 0x24, 0x30)
ROW_ALT  = RGBColor(0x12, 0x16, 0x1E)
DIVIDER  = RGBColor(0x2A, 0x30, 0x3E)
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = color
    run.font.name      = "Segoe UI"
    return txb


def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 0.06, ACCENT)
    add_text(slide, title, 0.5, 0.16, 10, 0.6, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.7, 12, 0.4, size=12, color=MUTED)
    add_rect(slide, 0.5, 1.04, 12.33, 0.02, SURFACE)


def card(slide, x, y, w, h, fill=SURFACE):
    add_rect(slide, x, y, w, h, fill)


def severity_pill(slide, sev, x, y):
    c = {"critical": CRITICAL, "warning": WARNING, "info": INFO_COL}.get(sev, MUTED)
    add_rect(slide, x, y, 1.3, 0.28, c)
    add_text(slide, sev.upper(), x, y + 0.02, 1.3, 0.26,
             size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 13.33, 0.08, ACCENT)
    add_rect(s, 0, 7.42, 13.33, 0.08, ACCENT)
    add_rect(s, 0.6, 1.6, 0.85, 0.85, ACCENT)
    add_text(s, "SL", 0.6, 1.64, 0.85, 0.77, size=28, bold=True, color=BG, align=PP_ALIGN.CENTER)
    add_text(s, "Servallab", 1.6, 1.55, 9, 0.72, size=48, bold=True, color=WHITE)
    add_text(s, "QC Engine", 1.6, 2.25, 9, 0.55, size=28, color=ACCENT)
    add_text(s, "Technical Architecture & Check Reference", 1.6, 2.9, 10, 0.5, size=16, color=MUTED)
    add_rect(s, 0.6, 3.55, 8, 0.02, SURFACE)
    add_text(s, "CATI Survey Data Quality Control", 0.6, 3.72, 9, 0.4, size=13, color=MUTED, italic=True)
    add_text(s, "Stack: FastAPI · React · Python · Groq LLM", 0.6, 4.12, 9, 0.4, size=13, color=MUTED)


def slide_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "System Architecture", "Two deployment surfaces sharing the same Python QC core")

    card(s, 0.5, 1.15, 5.8, 5.8, SURFACE)
    add_text(s, "Production Stack", 0.75, 1.28, 5.2, 0.4, size=14, bold=True, color=ACCENT)

    fe_blocks = [
        ("React + Vite + TypeScript", "Browser UI — tab-based dashboard"),
        ("Tailwind CSS + Zustand",    "Styling and in-memory app state"),
        ("FastAPI (Python)",          "REST API on port 8000"),
        ("Uvicorn ASGI",             "Async server, auto-reload in dev"),
        ("Axios",                    "Frontend → backend HTTP wrappers"),
    ]
    for i, (t, d) in enumerate(fe_blocks):
        y = 1.78 + i * 0.9
        add_rect(s, 0.65, y, 5.4, 0.76, DARKISH)
        add_text(s, t, 0.85, y + 0.04, 5.0, 0.35, size=12, bold=True, color=WHITE)
        add_text(s, d, 0.85, y + 0.39, 5.0, 0.3, size=10, color=MUTED)

    add_text(s, "→", 6.2, 3.6, 0.5, 0.5, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    card(s, 6.75, 1.15, 6.05, 5.8, SURFACE)
    add_text(s, "Python QC Core", 7.0, 1.28, 5.5, 0.4, size=14, bold=True, color=ACCENT2)

    core_blocks = [
        ("core/loader.py",       "Ingests CSV / XLSX / SAV → DataFrame"),
        ("core/cleaner.py",      "Normalises nulls, coerces types"),
        ("core/rule_engine.py",  "Reads config, instantiates checks, runs pipeline"),
        ("checks/ modules",      "11 check classes extending BaseCheck"),
        ("core/reporter.py",     "Excel / CSV / JSON output"),
    ]
    for i, (t, d) in enumerate(core_blocks):
        y = 1.78 + i * 0.9
        add_rect(s, 6.9, y, 5.7, 0.76, DARKISH)
        add_text(s, t, 7.1, y + 0.04, 5.3, 0.35, size=12, bold=True, color=WHITE)
        add_text(s, d, 7.1, y + 0.39, 5.3, 0.3, size=10, color=MUTED)


def slide_tools(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Tools & Why They Were Chosen", "Each tool selected for a specific technical need")

    tools = [
        ("FastAPI", "REST API framework", ACCENT,
         "Automatic OpenAPI docs, async support, Pydantic validation. Runs QC jobs in background threads without blocking the HTTP layer."),
        ("React + Vite", "Frontend UI", ACCENT2,
         "Component-based tab dashboard. Vite gives sub-second HMR in dev. TypeScript catches config shape mismatches at compile time."),
        ("Zustand", "State management", RGBColor(0xA0, 0x78, 0xF0),
         "Lightweight alternative to Redux. Single store holds file state, job polling, config, and UI flags — no boilerplate providers."),
        ("Pandas + NumPy", "Data processing", WARNING,
         "Industry-standard for tabular data. NumPy vectorised loops run straightlining scoring 3–5× faster than pandas .apply()."),
        ("Groq API", "Verbatim AI scoring", ACCENT,
         "Free-tier LLM inference fast enough to score open-ended responses in batches. Used over OpenAI to keep costs near zero."),
        ("python-pptx", "Report generation", MUTED,
         "Programmatic Excel and PowerPoint output without Microsoft Office on the server."),
    ]

    for i, (name, role, color, reason) in enumerate(tools):
        col, row = i % 3, i // 3
        x = 0.5 + col * 4.27
        y = 1.18 + row * 2.88
        card(s, x, y, 4.0, 2.68, SURFACE)
        add_rect(s, x, y, 4.0, 0.06, color)
        add_text(s, name, x + 0.15, y + 0.12, 3.7, 0.38, size=15, bold=True, color=WHITE)
        add_text(s, role, x + 0.15, y + 0.5,  3.7, 0.28, size=10, color=color)
        add_text(s, reason, x + 0.15, y + 0.85, 3.7, 1.62, size=10, color=MUTED)


def slide_pipeline(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "QC Pipeline", "Each uploaded file passes through four deterministic stages")

    stages = [
        ("1  Load", "core/loader.py",
         ["Accepts CSV, XLSX, SAV", "Reads into pandas DataFrame", "Reports row/column counts"]),
        ("2  Clean", "core/cleaner.py",
         ["Normalises null representations", "Coerces column types", "Strips whitespace"]),
        ("3  Validate", "rule_engine.py + checks/",
         ["RuleEngine reads JSON config", "Instantiates enabled check classes",
          "Calls .run(df) on each", "Collects CheckResult objects"]),
        ("4  Report", "reporter.py / routers/qc.py",
         ["API: serialises to JSON", "CLI: writes Excel + CSV",
          "Severity breakdown returned", "Flagged rows with _cols preserved"]),
    ]

    for i, (title, module, bullets) in enumerate(stages):
        x = 0.4 + i * 3.22
        card(s, x, 1.18, 3.0, 5.6, SURFACE)
        add_rect(s, x, 1.18, 3.0, 0.06, ACCENT if i % 2 == 0 else ACCENT2)
        add_text(s, title, x + 0.15, 1.3, 2.7, 0.5, size=16, bold=True, color=WHITE)
        add_text(s, module, x + 0.15, 1.78, 2.7, 0.35, size=9, color=ACCENT if i % 2 == 0 else ACCENT2, italic=True)
        for j, b in enumerate(bullets):
            add_text(s, "▸  " + b, x + 0.15, 2.28 + j * 0.68, 2.7, 0.62, size=11, color=MUTED)
        if i < 3:
            add_text(s, "→", x + 3.02, 3.8, 0.22, 0.4, size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_check_overview(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Check Inventory", "All check classes — module, severity, trigger condition")

    checks = [
        ("Missing Value Check",          "missing_checks.py",  "warning",  "Any required column is null"),
        ("High Missing Column Check",    "missing_checks.py",  "info",     "Column null rate exceeds threshold"),
        ("Range Check",                  "range_checks.py",    "warning",  "Numeric value outside [min, max]"),
        ("Duration Check",               "range_checks.py",    "warning",  "Interview time below min or above max"),
        ("Pattern Check (Regex)",        "pattern_checks.py",  "warning",  "Value doesn't match required regex"),
        ("Anomaly Check (IQR)",          "pattern_checks.py",  "info",     "Statistical outlier via IQR fence"),
        ("Logic / Skip-Logic Check",     "logic_checks.py",    "critical", "IF condition met but THEN violated"),
        ("Duplicate Check",              "logic_checks.py",    "critical", "Row duplicated on key columns"),
        ("Straightlining Check",         "advanced_checks.py", "warning",  "≥ threshold% same answer across grid"),
        ("Interviewer Duration Anomaly", "advanced_checks.py", "warning",  "Interviewer avg duration outside IQR"),
        ("Interviewer Productivity",     "advanced_checks.py", "warning",  "Interview count outside IQR range"),
        ("Consent / Eligibility",        "advanced_checks.py", "critical", "Ineligible respondent has answers"),
        ("Fabrication Detection",        "advanced_checks.py", "critical", "Sequential IDs or low numeric variance"),
        ("Verbatim Quality (AI)",        "verbatim_checks.py", "warning",  "Open-end dimension score < min_score"),
    ]

    add_rect(s, 0.4, 1.1, 12.5, 0.34, DARKISH)
    for label, x in [("Check", 0.55), ("Module", 4.6), ("Severity", 7.5), ("Trigger", 9.2)]:
        add_text(s, label, x, 1.14, 2.0, 0.28, size=10, bold=True, color=MUTED)

    for i, (name, module, sev, trigger) in enumerate(checks):
        y = 1.48 + i * 0.41
        add_rect(s, 0.4, y, 12.5, 0.41, SURFACE if i % 2 == 0 else ROW_ALT)
        add_text(s, name,   0.55, y + 0.06, 4.0, 0.3, size=10, color=WHITE)
        add_text(s, module, 4.6,  y + 0.06, 2.8, 0.3, size=9,  color=MUTED, italic=True)
        sev_c = {"critical": CRITICAL, "warning": WARNING, "info": INFO_COL}.get(sev, MUTED)
        add_text(s, sev.upper(), 7.5, y + 0.06, 1.6, 0.3, size=9, bold=True, color=sev_c)
        add_text(s, trigger, 9.2, y + 0.06, 3.6, 0.3, size=9, color=MUTED)


def slide_severity(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Severity Framework", "Every check result carries one of three severity levels")

    levels = [
        ("CRITICAL", CRITICAL,
         "Data integrity is compromised. These rows must be reviewed before sign-off.",
         ["Logic violations (IF/THEN broken)", "Duplicate respondent records",
          "Consent / eligibility breach", "Fabrication indicators"]),
        ("WARNING", WARNING,
         "Suspicious pattern that may indicate poor data quality or interviewer error.",
         ["Straightlining (satisficing)", "Duration anomaly (too fast / slow)",
          "Interviewer productivity outlier", "Out-of-range numeric values",
          "Regex format violations", "Low verbatim quality score"]),
        ("INFO", INFO_COL,
         "Informational — no immediate action required but worth monitoring.",
         ["High column missing rate (structural issue)",
          "Statistical anomaly detected by IQR"]),
    ]

    for i, (label, color, desc, bullets) in enumerate(levels):
        x = 0.5 + i * 4.27
        card(s, x, 1.18, 4.0, 5.8, SURFACE)
        add_rect(s, x, 1.18, 4.0, 0.08, color)
        add_text(s, label, x + 0.15, 1.32, 3.7, 0.45, size=20, bold=True, color=WHITE)
        add_text(s, desc,  x + 0.15, 1.82, 3.7, 0.85, size=10, color=MUTED)
        add_rect(s, x + 0.15, 2.72, 3.7, 0.02, DIVIDER)
        for j, b in enumerate(bullets):
            add_text(s, "▸  " + b, x + 0.15, 2.84 + j * 0.62, 3.7, 0.58, size=10, color=WHITE)


def slide_missing(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Missing Value Checks", "missing_checks.py  ·  MissingValueCheck + HighMissingColumnCheck")

    card(s, 0.5, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 5.9, 0.06, WARNING)
    add_text(s, "MissingValueCheck", 0.7, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)
    severity_pill(s, "warning", 0.7, 1.76)

    r = [
        ("What it does", "Flags rows where any required column is null."),
        ("With threshold", "Identifies columns whose null rate > threshold, then flags rows with nulls in those columns."),
        ("Without threshold", "Flags every row that has a null in any of the checked columns."),
        ("Output columns", "_missing_columns — list of which columns were null on that row."),
        ("Config key", "missing_threshold: 0.1  (10% → flag the column)"),
    ]
    y = 2.18
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=ACCENT)
        add_text(s, val,   2.85, y, 3.4, 0.55, size=10, color=MUTED)
        y += 0.66

    card(s, 6.8, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 6.8, 1.15, 5.9, 0.06, INFO_COL)
    add_text(s, "HighMissingColumnCheck", 7.0, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)
    severity_pill(s, "info", 7.0, 1.76)

    r2 = [
        ("What it does", "Scans every column and reports those whose null rate exceeds the threshold. Column-level summary, not per-row."),
        ("Output", "One summary row per bad column:\n  column name + missing_rate"),
        ("Use case", "Detects structural collection failures — e.g. a field never recorded for a wave."),
        ("Config key", "missing_threshold: 0.2  (20% default)"),
    ]
    y = 2.18
    for label, val in r2:
        add_text(s, label, 7.0, y, 2.1, 0.35, size=10, bold=True, color=INFO_COL)
        add_text(s, val,   9.15, y, 3.3, 0.68, size=10, color=MUTED)
        y += 0.82


def slide_range_duration(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Range & Duration Checks", "range_checks.py  ·  RangeCheck + DurationCheck")

    card(s, 0.5, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 5.9, 0.06, WARNING)
    add_text(s, "RangeCheck", 0.7, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)
    severity_pill(s, "warning", 0.7, 1.76)

    r = [
        ("What it does", "For each rule, casts column to numeric and flags rows where value < min or > max."),
        ("Scoring", "Binary: in-range = pass, out-of-range = flagged. No partial score."),
        ("Output columns", "_range_issue (column name)\n_value (actual value)\n_expected_range [min, max]"),
        ("Multiple rules", "Each rule processed independently; all flags concatenated."),
        ("Config", '{"column":"age","min":18,"max":99}'),
    ]
    y = 2.18
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=ACCENT)
        add_text(s, val,   2.85, y, 3.4, 0.62, size=10, color=MUTED)
        y += 0.76

    card(s, 6.8, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 6.8, 1.15, 5.9, 0.06, WARNING)
    add_text(s, "DurationCheck", 7.0, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)
    severity_pill(s, "warning", 7.0, 1.76)

    r2 = [
        ("What it does", "Flags interviews whose duration (minutes) falls below min or above max."),
        ("Scoring", "< min_minutes → too_short\n> max_minutes → too_long\nBoth categories counted separately in metadata."),
        ("Output columns", "_duration_issue (the raw duration value)"),
        ("Metadata", "min_expected, max_expected\ntoo_short count, too_long count"),
        ("Config", "column: duration_minutes\nmin_expected: 5,  max_expected: 120"),
    ]
    y = 2.18
    for label, val in r2:
        add_text(s, label, 7.0, y, 2.1, 0.35, size=10, bold=True, color=ACCENT)
        add_text(s, val,   9.15, y, 3.3, 0.62, size=10, color=MUTED)
        y += 0.76


def slide_pattern_anomaly(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Pattern & Anomaly Checks", "pattern_checks.py  ·  PatternCheck + AnomalyCheck")

    card(s, 0.5, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 5.9, 0.06, WARNING)
    add_text(s, "PatternCheck (Regex)", 0.7, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)
    severity_pill(s, "warning", 0.7, 1.76)

    r = [
        ("What it does", "Compiles each rule's regex and flags non-null values that do not match."),
        ("Scoring", "Binary match/fail. Null values are skipped (not flagged)."),
        ("Output columns", "_pattern_issue (description)\n_invalid_value (the actual bad value)"),
        ("Use cases", "Phone formats, ID codes, date strings, postal codes, fixed-choice codes."),
        ("Config", '{"column":"phone",\n "pattern":"^\\d{10}$",\n "description":"10-digit phone"}'),
    ]
    y = 2.18
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=ACCENT)
        add_text(s, val,   2.85, y, 3.4, 0.68, size=10, color=MUTED)
        y += 0.76

    card(s, 6.8, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 6.8, 1.15, 5.9, 0.06, INFO_COL)
    add_text(s, "AnomalyCheck (IQR)", 7.0, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)
    severity_pill(s, "info", 7.0, 1.76)

    r2 = [
        ("What it does", "Flags statistical outliers in numeric columns using the IQR fence method."),
        ("Method", "Q1 − k×IQR  to  Q3 + k×IQR\nwhere k = multiplier (default 1.5).\nValues outside the fence are flagged."),
        ("Multiplier guide", "1.5 = standard outlier\n3.0 = extreme outlier only"),
        ("Output columns", "_anomaly_column, _anomaly_value\n_anomaly_bounds [lower, upper]"),
        ("Note", "Severity is INFO — exploratory signal, not a definitive flag."),
    ]
    y = 2.18
    for label, val in r2:
        add_text(s, label, 7.0, y, 2.1, 0.35, size=10, bold=True, color=INFO_COL)
        add_text(s, val,   9.15, y, 3.3, 0.68, size=10, color=MUTED)
        y += 0.76


def slide_logic(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Logic / Skip-Logic Check", "logic_checks.py  ·  LogicCheck  ·  severity: CRITICAL")

    card(s, 0.5, 1.15, 6.2, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 6.2, 0.06, CRITICAL)

    r = [
        ("What it does",
         "Validates user-defined IF→THEN rules. Each rule fires when ALL if_conditions are met, "
         "then checks that EACH then_condition also holds. A then_condition that fails = violation."),
        ("Condition logic",
         "Multiple if_conditions are AND-ed.\nEach then_condition is evaluated independently — "
         "a row can violate one or several."),
        ("Output columns",
         "_logic_rule (rule description)\n_violated_column (which THEN column failed)\n"
         "_violated_condition (the exact condition string)"),
        ("Backwards compat",
         "Supports legacy flat format (if_column / if_value / then_column / then_condition) "
         "so old Python configs still load without change."),
        ("Severity", "CRITICAL — logic violations indicate data collection errors or interviewer manipulation."),
    ]
    y = 1.35
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=CRITICAL)
        add_text(s, val,   2.85, y, 3.65, 0.85, size=10, color=MUTED)
        y += 0.96

    card(s, 7.1, 1.15, 5.7, 2.6, SURFACE)
    add_text(s, "Supported Operators", 7.3, 1.22, 5.2, 0.38, size=12, bold=True, color=WHITE)
    ops = [
        (">  <  >=  <=  ==  !=", "Numeric and string comparisons"),
        ("is_null  /  not_null",  "Presence / absence check"),
        ("is_numeric  /  is_string", "Type assertion"),
        ("in_list  /  not_in_list", "Membership (value = list)"),
    ]
    for i, (op, desc) in enumerate(ops):
        y = 1.68 + i * 0.46
        add_text(s, op,   7.3,  y, 2.5, 0.36, size=10, bold=True, color=ACCENT)
        add_text(s, desc, 9.85, y, 2.7, 0.36, size=10, color=MUTED)

    card(s, 7.1, 3.88, 5.7, 3.12, DARKER)
    add_text(s, "Example Rule", 7.3, 3.96, 5.2, 0.34, size=11, bold=True, color=MUTED)
    example = (
        '{\n'
        '  "description": "Under-18: no salary",\n'
        '  "if_conditions": [\n'
        '    {"column":"age","operator":"<","value":18}\n'
        '  ],\n'
        '  "then_conditions": [\n'
        '    {"column":"salary","operator":"is_null"}\n'
        '  ]\n'
        '}'
    )
    add_text(s, example, 7.3, 4.38, 5.3, 2.4, size=9, color=ACCENT)


def slide_duplicates(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Duplicate Check", "logic_checks.py  ·  DuplicateCheck  ·  severity: CRITICAL")

    card(s, 0.5, 1.15, 12.3, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 12.3, 0.06, CRITICAL)

    r = [
        ("What it does",
         "Identifies records that are identical across a set of key columns. Uses all columns if no subset is specified."),
        ("Detection method",
         "Casts all key columns to string, concatenates with '|', then marks every row where the key appears more than once "
         "(keep=False — ALL copies flagged, not just the repeat)."),
        ("Subset columns",
         "If subset_columns is empty the check uses all columns. Typically set to the respondent ID column, "
         "or a combination of ID + phone."),
        ("Output columns", "_dupe_key — the list of columns used as the duplicate key."),
        ("Why CRITICAL",
         "Duplicate respondents inflate sample sizes, skew weighted estimates, and often indicate ghost interviews "
         "entered by interviewers to meet quotas."),
        ("Config",
         'duplicate_check: {"enabled": true, "subset_columns": ["respondent_id"]}'),
    ]
    y = 1.38
    for label, val in r:
        add_text(s, label, 0.7, y, 2.2, 0.35, size=11, bold=True, color=CRITICAL)
        add_text(s, val,   3.0,  y, 9.5, 0.72, size=11, color=MUTED)
        y += 0.85


def slide_straightlining(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Straightlining Check", "advanced_checks.py  ·  StraightliningCheck  ·  severity: WARNING")

    card(s, 0.5, 1.15, 6.2, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 6.2, 0.06, WARNING)

    r = [
        ("What it does",
         "Detects respondents who gave the same answer across a grid of scale questions — a sign of satisficing or fabrication."),
        ("Algorithm",
         "Per row: count non-null answers, find modal value, compute score = modal_count / n_answers. "
         "Vectorised via NumPy for 3–5× speedup over .apply()."),
        ("Threshold",
         "score ≥ threshold → flagged.\n0.9 = ≥90% answers identical.\n1.0 = flags only perfect straightliners."),
        ("Min questions",
         "Rows with fewer than min_questions non-null answers are skipped — insufficient data to judge."),
        ("Interviewer view",
         "When interviewer_column is set, produces per-interviewer summary: straightliner count and % rate."),
        ("Output columns",
         "_sl_score (0–1)\n_sl_modal_answer (the repeated value)"),
    ]
    y = 1.32
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=WARNING)
        add_text(s, val,   2.85, y, 3.65, 0.82, size=10, color=MUTED)
        y += 0.9

    card(s, 7.1, 1.15, 5.7, 5.85, SURFACE)
    add_text(s, "Scoring Criteria — _sl_score", 7.3, 1.22, 5.2, 0.42, size=13, bold=True, color=WHITE)

    score_rows = [
        ("1.00",       "Perfect straightliner — every answer identical",        CRITICAL),
        ("0.90–0.99",  "Near-perfect — 90%+ same (default flag zone)",          WARNING),
        ("0.80–0.89",  "Suspicious — flag with stricter threshold setting",     WARNING),
        ("< 0.80",     "Normal variation — not flagged",                        ACCENT),
    ]
    y = 1.74
    for score, label, color in score_rows:
        add_rect(s, 7.15, y, 5.6, 0.65, DARK)
        add_text(s, score, 7.3, y + 0.12, 1.5, 0.38, size=13, bold=True, color=color)
        add_text(s, label, 8.85, y + 0.12, 3.8, 0.38, size=10, color=MUTED)
        y += 0.72

    add_rect(s, 7.15, y + 0.1, 5.6, 0.02, DIVIDER)
    add_text(s, "Recommended Settings", 7.3, y + 0.22, 5.2, 0.35, size=11, bold=True, color=WHITE)
    for j, line in enumerate([
        "threshold: 0.9   (90%+ same answer)",
        "min_questions: 3   (minimum grid columns)",
        "Lower to 0.8 for short grids (4-item scales)",
    ]):
        add_text(s, "▸  " + line, 7.3, y + 0.65 + j * 0.44, 5.2, 0.38, size=10, color=MUTED)


def slide_interviewer_duration(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Interviewer Duration Anomaly", "advanced_checks.py  ·  InterviewerDurationCheck  ·  severity: WARNING")

    card(s, 0.5, 1.15, 6.2, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 6.2, 0.06, WARNING)

    r = [
        ("What it does",
         "Computes each interviewer's mean interview duration, then flags those whose mean is a statistical outlier vs the peer group."),
        ("IQR method",
         "lower = Q1 − k×IQR\nupper = Q3 + k×IQR\nwhere k = multiplier (default 1.5)."),
        ("min_interviews",
         "Interviewers with fewer than min_interviews records are excluded from the IQR calculation."),
        ("Min group",
         "Requires ≥ 4 interviewers with enough data. Returns empty result otherwise."),
        ("Output columns",
         "mean_duration, interview_count\n_duration_issue: 'too_fast' or 'too_slow'\n_expected_range [lo, hi] mins"),
        ("Flagged scope",
         "ALL interviews by outlier interviewers are flagged — not just averages."),
    ]
    y = 1.32
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=WARNING)
        add_text(s, val,   2.85, y, 3.65, 0.85, size=10, color=MUTED)
        y += 0.9

    card(s, 7.1, 1.15, 5.7, 5.85, SURFACE)
    add_text(s, "IQR Multiplier Guide", 7.3, 1.22, 5.2, 0.42, size=13, bold=True, color=WHITE)

    for mult, desc, color in [
        ("k = 1.5", "Standard fence — catches moderately fast/slow interviewers", WARNING),
        ("k = 2.0", "Conservative — obvious outliers only (fewer false positives)", INFO_COL),
        ("k = 3.0", "Extreme outlier only — use when duration naturally varies a lot", ACCENT),
    ]:
        add_rect(s, 7.15, y := (y if 'y' not in dir() else y), 5.6, 0.72, DARK)
        add_text(s, mult, 7.3, y + 0.14, 1.4, 0.38, size=13, bold=True, color=color)
        add_text(s, desc, 8.75, y + 0.14, 3.8, 0.45, size=10, color=MUTED)

    mults = [
        ("k = 1.5", "Standard fence — catches moderately fast/slow interviewers", WARNING),
        ("k = 2.0", "Conservative — obvious outliers only (fewer false positives)", INFO_COL),
        ("k = 3.0", "Extreme outlier only — use when duration naturally varies a lot", ACCENT),
    ]
    y = 1.74
    for mult, desc, color in mults:
        add_rect(s, 7.15, y, 5.6, 0.72, DARK)
        add_text(s, mult, 7.3, y + 0.14, 1.4, 0.38, size=13, bold=True, color=color)
        add_text(s, desc, 8.75, y + 0.14, 3.8, 0.45, size=10, color=MUTED)
        y += 0.82

    add_text(s, "Flag labels", 7.3, y + 0.2, 5.2, 0.35, size=11, bold=True, color=WHITE)
    for i, (lbl, desc, color) in enumerate([
        ("too_fast", "Mean duration < lower fence — rushed / fabricated?", CRITICAL),
        ("too_slow", "Mean duration > upper fence — distracted / probing?", WARNING),
    ]):
        add_rect(s, 7.15, y + 0.65 + i * 0.72, 5.6, 0.65, DARK)
        add_text(s, lbl,  7.3,  y + 0.74 + i * 0.72, 1.5, 0.38, size=11, bold=True, color=color)
        add_text(s, desc, 8.85, y + 0.74 + i * 0.72, 3.8, 0.38, size=10, color=MUTED)


def slide_productivity(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Interviewer Productivity Check", "advanced_checks.py  ·  InterviewerProductivityCheck  ·  severity: WARNING")

    card(s, 0.5, 1.15, 6.2, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 6.2, 0.06, WARNING)

    r = [
        ("What it does",
         "Counts completed interviews per interviewer, then uses IQR to identify those with suspiciously many or few."),
        ("IQR method",
         "lower = max(0, Q1 − k×IQR)\nupper = Q3 + k×IQR"),
        ("unusually_high",
         "Count > upper fence. Primary fraud signal — fabricated records to hit targets. "
         "Use with duration check for stronger evidence."),
        ("unusually_low",
         "Count < lower fence. May indicate laziness, early dropout, or data recording failure."),
        ("Output columns",
         "interview_count\n_productivity_issue: 'unusually_high' or 'unusually_low'\n_expected_range [lo, hi]"),
        ("Config", "interviewer_column: 'interviewer_id'\nmultiplier: 1.5"),
    ]
    y = 1.32
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=WARNING)
        add_text(s, val,   2.85, y, 3.65, 0.85, size=10, color=MUTED)
        y += 0.9

    card(s, 7.1, 1.15, 5.7, 5.85, SURFACE)
    add_text(s, "Combined Signal", 7.3, 1.22, 5.2, 0.42, size=13, bold=True, color=WHITE)
    add_text(s,
        "Productivity and duration checks are designed to be used together. "
        "A strong fabrication signal requires both:\n\n"
        "▸  High interview count  (productivity outlier)\n"
        "▸  Very short average duration  (duration: too_fast)\n\n"
        "Either flag alone may have an innocent explanation. "
        "Both together strongly suggest ghost interviews.",
        7.3, 1.72, 5.3, 2.5, size=11, color=MUTED)

    add_rect(s, 7.15, 4.35, 5.6, 0.02, DIVIDER)
    add_text(s, "Scoring Criteria", 7.3, 4.48, 5.2, 0.35, size=11, bold=True, color=WHITE)

    for i, (lbl, desc, color) in enumerate([
        ("unusually_high", "> Q3 + k×IQR interviews", CRITICAL),
        ("unusually_low",  "< Q1 − k×IQR interviews", INFO_COL),
        ("normal",         "Within the IQR fence — not flagged", ACCENT),
    ]):
        add_rect(s, 7.15, 4.92 + i * 0.62, 5.6, 0.56, DARK)
        add_text(s, lbl,  7.3,  4.99 + i * 0.62, 1.8, 0.38, size=10, bold=True, color=color)
        add_text(s, desc, 9.15, 4.99 + i * 0.62, 3.5, 0.38, size=10, color=MUTED)


def slide_consent(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Consent / Eligibility Check", "advanced_checks.py  ·  ConsentEligibilityCheck  ·  severity: CRITICAL")

    card(s, 0.5, 1.15, 12.3, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 12.3, 0.06, CRITICAL)

    r = [
        ("What it does",
         "Identifies respondents disqualified by a screener question who still have data in subsequent questions — "
         "the interview should not have continued."),
        ("Detection logic",
         "Step 1: Find all rows where screener_column meets the disqualify condition (same operator engine as LogicCheck).\n"
         "Step 2: Of those rows, flag any that have at least one non-null value in any of the subsequent_columns list."),
        ("Output columns",
         "_screener_value (what the screener column contained)\n"
         "_populated_columns (list of subsequent columns that had data despite disqualification)"),
        ("Severity",
         "CRITICAL — a disqualified respondent with subsequent data means the screener was ignored, "
         "or data was fabricated for an ineligible person."),
        ("Operators",
         "Uses the full _evaluate_condition() engine — supports ==, !=, in_list, not_in_list, is_null, not_null, >, <, >=, <=."),
        ("Config",
         'screener_column: "eligible"  ·  disqualify_operator: "!="  ·  disqualify_value: "yes"\n'
         'subsequent_columns: ["q1", "q2", "q3"]'),
    ]
    y = 1.38
    for label, val in r:
        add_text(s, label, 0.7, y, 2.2, 0.35, size=11, bold=True, color=CRITICAL)
        add_text(s, val,   3.0,  y, 9.5, 0.78, size=11, color=MUTED)
        y += 0.87


def slide_fabrication(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Fabrication Detection", "advanced_checks.py  ·  FabricationCheck  ·  severity: CRITICAL")

    card(s, 0.5, 1.15, 6.2, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 6.2, 0.06, CRITICAL)
    add_text(s, "Detection A — Sequential IDs", 0.7, 1.28, 5.7, 0.42, size=13, bold=True, color=WHITE)

    a_rows = [
        ("What it is", "Consecutive numeric respondent IDs in a random sample are suspicious — real field collection rarely produces perfectly sequential IDs."),
        ("Algorithm",  "Sorts IDs numerically. Walks the sequence looking for runs where each value = previous + 1. Flags any run ≥ sequence_run_length."),
        ("Config",     "id_column: 'respondent_id'\nsequence_run_length: 5"),
        ("Output",     "_fabrication_type = 'sequential_ids'"),
    ]
    y = 1.8
    for label, val in a_rows:
        add_text(s, label, 0.7, y, 2.0, 0.35, size=10, bold=True, color=CRITICAL)
        add_text(s, val,   2.75, y, 3.75, 0.65, size=10, color=MUTED)
        y += 0.84

    card(s, 7.1, 1.15, 5.7, 5.85, SURFACE)
    add_rect(s, 7.1, 1.15, 5.7, 0.06, CRITICAL)
    add_text(s, "Detection B — Low Variance", 7.3, 1.28, 5.2, 0.42, size=13, bold=True, color=WHITE)

    b_rows = [
        ("What it is",  "If an interviewer's respondents all gave nearly identical numeric answers, the data may have been copied or invented."),
        ("Algorithm",   "For each numeric column, compute each interviewer's std dev. "
                        "If interviewer_std / global_std < variance_threshold → suspicious."),
        ("Threshold",   "variance_threshold: 0.1 means the interviewer's SD must be < 10% of the global SD to be flagged."),
        ("Config",      "numeric_columns: ['q1','q2','q3']\nvariance_threshold: 0.1"),
        ("Output",      "_fabrication_type = 'low_variance_{col}'\n_column_checked = which column triggered"),
    ]
    y = 1.8
    for label, val in b_rows:
        add_text(s, label, 7.3, y, 1.85, 0.35, size=10, bold=True, color=CRITICAL)
        add_text(s, val,   9.2,  y, 3.45, 0.68, size=10, color=MUTED)
        y += 0.84


def slide_verbatim(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Verbatim Quality Check (AI)", "verbatim_checks.py  ·  VerbatimQualityCheck  ·  severity: WARNING")

    card(s, 0.5, 1.15, 6.2, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 6.2, 0.06, ACCENT)

    r = [
        ("What it does",
         "Sends open-ended survey responses to a Groq-hosted LLM and asks it to score each on four quality dimensions."),
        ("Batching",
         "Responses grouped in batches of 10 — one API call per batch. 100 responses = 10 API calls."),
        ("Sampling",
         "If dataset exceeds sample_size, a random sample is scored (random_state=42). Full dataset not required."),
        ("API key order",
         "1. GROQ_API_KEY env var (server key)\n2. User's personal key from Settings\n"
         "On HTTP 429 rate-limit, retries with personal key before giving up."),
        ("Flag condition",
         "Flagged if ANY dimension scores < min_score, OR if any boolean flag (gibberish / copy_paste / too_short) is true."),
        ("Models",
         "llama-3.1-8b-instant (default, fastest)\nllama-3.3-70b-versatile\ngemma2-9b-it"),
    ]
    y = 1.32
    for label, val in r:
        add_text(s, label, 0.7, y, 2.1, 0.35, size=10, bold=True, color=ACCENT)
        add_text(s, val,   2.85, y, 3.65, 0.82, size=10, color=MUTED)
        y += 0.9

    card(s, 7.1, 1.15, 5.7, 5.85, SURFACE)
    add_text(s, "AI Scoring Dimensions (1–5)", 7.3, 1.22, 5.2, 0.42, size=13, bold=True, color=WHITE)

    for i, (dim, desc) in enumerate([
        ("grammar",        "Grammatical correctness of the response"),
        ("coherence",      "Whether the response makes logical sense"),
        ("relevance",      "Whether it answers the question asked"),
        ("length_quality", "Appropriate length — not too short or padded"),
    ]):
        y = 1.74 + i * 0.66
        add_rect(s, 7.15, y, 5.6, 0.6, DARK)
        add_text(s, dim,  7.3,  y + 0.1, 1.9, 0.38, size=11, bold=True, color=ACCENT)
        add_text(s, desc, 9.25, y + 0.1, 3.35, 0.38, size=10, color=MUTED)

    y = 4.44
    add_rect(s, 7.15, y, 5.6, 0.02, DIVIDER)
    add_text(s, "Boolean Flags (any = flagged)", 7.3, y + 0.12, 5.2, 0.35, size=11, bold=True, color=WHITE)

    for i, (flag, desc) in enumerate([
        ("gibberish",  "Random characters or nonsense text"),
        ("copy_paste", "Identical / near-identical to another response"),
        ("too_short",  "Insufficiently detailed for the question"),
    ]):
        add_rect(s, 7.15, y + 0.56 + i * 0.55, 5.6, 0.5, DARK)
        add_text(s, flag, 7.3,  y + 0.63 + i * 0.55, 1.5, 0.38, size=10, bold=True, color=CRITICAL)
        add_text(s, desc, 8.85, y + 0.63 + i * 0.55, 3.8, 0.38, size=10, color=MUTED)


def slide_verbatim_scoring(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Verbatim Scoring Criteria — Detailed", "How the LLM grades each dimension on a 1–5 scale")

    dims = [
        ("grammar", ACCENT,
         [("5", "Perfect grammar, no errors"),
          ("4", "Minor errors that don't impede understanding"),
          ("3", "Noticeable errors but message is clear"),
          ("2", "Significant errors, hard to parse"),
          ("1", "Unintelligible or completely ungrammatical")]),
        ("coherence", ACCENT2,
         [("5", "Fully coherent, logical flow"),
          ("4", "Mostly coherent, minor inconsistency"),
          ("3", "Some confusion but overall meaning survives"),
          ("2", "Disjointed or contradictory"),
          ("1", "No discernible logical structure")]),
        ("relevance", WARNING,
         [("5", "Directly answers the question"),
          ("4", "Mostly relevant, minor tangent"),
          ("3", "Partially relevant — touches the topic"),
          ("2", "Mostly off-topic"),
          ("1", "Completely irrelevant or refuses")]),
        ("length_quality", RGBColor(0xA0, 0x78, 0xF0),
         [("5", "Appropriate, substantive length"),
          ("4", "Slightly short/long but acceptable"),
          ("3", "Noticeably too brief or repetitively padded"),
          ("2", "Very short (1–3 words) or excessively repeated"),
          ("1", "Single character, placeholder, or blank-equivalent")]),
    ]

    for i, (name, color, scores) in enumerate(dims):
        col, row = i % 2, i // 2
        x = 0.5 + col * 6.5
        y = 1.15 + row * 3.15

        card(s, x, y, 6.2, 3.0, SURFACE)
        add_rect(s, x, y, 6.2, 0.06, color)
        add_text(s, name, x + 0.2, y + 0.1, 5.7, 0.4, size=14, bold=True, color=WHITE)

        for j, (score, desc) in enumerate(scores):
            sy = y + 0.6 + j * 0.46
            sc = CRITICAL if score == "1" else WARNING if score == "2" else WARNING if score == "3" else ACCENT
            add_text(s, score, x + 0.2, sy, 0.4, 0.38, size=13, bold=True, color=sc)
            add_text(s, desc,  x + 0.65, sy, 5.4, 0.38, size=10, color=MUTED)

    add_text(s,
        "A response is flagged if ANY dimension < min_score (default 2) "
        "OR any boolean flag (gibberish / copy_paste / too_short) is true.",
        0.5, 7.1, 12.3, 0.32, size=11, color=MUTED, italic=True)


def slide_config_profiles(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Config Profiles", "Portable JSON files that define exactly which checks run and with what parameters")

    card(s, 0.5, 1.15, 6.2, 5.85, SURFACE)
    add_text(s, "Why Config Files?", 0.7, 1.22, 5.7, 0.42, size=13, bold=True, color=WHITE)

    r = [
        ("Repeatability",
         "The same config applied to Wave 1 data applies identically to Wave 2 — no manual re-entry of column names."),
        ("Portability",
         "A QC manager defines rules for a project once and shares the JSON file with field supervisors."),
        ("Auditability",
         "The config file is a record of what was checked and how — version-controllable alongside the data."),
        ("Template-driven",
         "The app ships a blank template with all sections pre-filled with example values "
         "and _info comment fields explaining each parameter."),
    ]
    y = 1.72
    for label, val in r:
        add_text(s, label, 0.7, y, 2.0, 0.35, size=10, bold=True, color=ACCENT)
        add_text(s, val,   2.75, y, 3.75, 0.7, size=10, color=MUTED)
        y += 0.88

    card(s, 7.1, 1.15, 5.7, 5.85, DARKER)
    add_text(s, "Template Structure", 7.3, 1.22, 5.2, 0.38, size=12, bold=True, color=MUTED)
    snippet = (
        '{\n'
        '  "_info": "Set enabled:true + fill column names",\n\n'
        '  "missing_threshold": 0.1,\n\n'
        '  "straightlining": {\n'
        '    "_info": "Flags same-answer grid respondents",\n'
        '    "enabled": false,\n'
        '    "question_columns": ["q1","q2","q3"],\n'
        '    "threshold": 0.9\n'
        '  },\n\n'
        '  "logic_rules": [\n'
        '    {\n'
        '      "description": "Consent → must answer Q1",\n'
        '      "if_conditions":   [...],\n'
        '      "then_conditions": [...]\n'
        '    }\n'
        '  ],\n\n'
        '  "fabrication_check": { "enabled": false, ... },\n'
        '  "verbatim_check":    { "enabled": false, ... }\n'
        '}'
    )
    add_text(s, snippet, 7.25, 1.68, 5.4, 5.1, size=9, color=ACCENT)


def slide_dashboard_overview(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Manager Dashboard", "Streamlit subdomain app — dashboard.servallab.com  ·  separate from QC engine")

    card(s, 0.5, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 5.9, 0.06, ACCENT2)
    add_text(s, "What It Is", 0.7, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)

    r = [
        ("Purpose",
         "A multi-project operations dashboard for QC managers and supervisors to track survey health "
         "across all active studies — not a replacement for the QC engine."),
        ("Stack",
         "Streamlit Python app. Deployed on Fly.io (dashboard.servallab.com). "
         "SQLite database via SQLAlchemy — persists project, interviewer, and QC history records."),
        ("Auth",
         "JWT-based login. Roles: admin, manager, viewer. "
         "Admin creates users; managers own projects; viewers can read but not write."),
        ("Relation to QC engine",
         "QC results saved from the React engine can be pushed here via 'Save to project'. "
         "The dashboard tracks trends over multiple waves, not individual file runs."),
        ("Default credentials",
         "First-run: admin@servallab.com / admin1234. Change immediately in production."),
    ]
    y = 1.75
    for label, val in r:
        add_text(s, label, 0.7, y, 2.0, 0.35, size=10, bold=True, color=ACCENT2)
        add_text(s, val,   2.75, y, 3.55, 0.78, size=10, color=MUTED)
        y += 0.9

    card(s, 6.8, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 6.8, 1.15, 5.9, 0.06, ACCENT2)
    add_text(s, "Dashboard Pages", 7.0, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)

    pages = [
        ("Project Overview",   "KPIs, wave history, open issues per project"),
        ("Quality Report",     "Flag rate trends by check type across waves"),
        ("Performance Report", "Interviewer-level flag rates and risk scores over time"),
        ("Backcheck Report",   "Backcheck completion rate and error tracking"),
        ("Listen-In Report",   "Listen-in session pass/fail tracking per supervisor"),
        ("Wave Comparison",    "Statistical shifts between survey waves"),
        ("Cancelled Interviews","Cancellations logged by reason and interviewer"),
        ("Admin",              "User management, project creation, role assignment"),
    ]
    y = 1.75
    for page, desc in pages:
        add_rect(s, 6.9, y, 5.6, 0.62, DARKISH)
        add_text(s, page, 7.05, y + 0.06, 2.2, 0.32, size=10, bold=True, color=WHITE)
        add_text(s, desc,  9.3, y + 0.06, 3.05, 0.5, size=9, color=MUTED)
        y += 0.7


def slide_dashboard_db(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Dashboard — Data Model", "SQLite (SQLAlchemy ORM)  ·  dashboard/database.py")

    tables = [
        ("User", ACCENT2,
         ["id, email, hashed_password", "role: admin / manager / viewer",
          "created_at, is_active"]),
        ("Project", ACCENT,
         ["id, name, status (active/closed)", "manager_id → User",
          "created_at, wave_count"]),
        ("UploadLog", WARNING,
         ["id, project_id → Project", "file_id, filename, wave_label",
          "total_records, flag_count", "uploaded_at"]),
        ("Interviewer", ACCENT2,
         ["id, interviewer_code (unique)", "name, supervisor_name",
          "created_at, is_active"]),
        ("InterviewerQuality", CRITICAL,
         ["interviewer_id → Interviewer", "project_id → Project",
          "total_interviews, duration_flags", "sl_flags, total_flags"]),
        ("BackcheckRecord", INFO_COL,
         ["interviewer_id → Interviewer", "project_id, bc_date",
          "total_errors, outcome"]),
        ("ListenInRecord", RGBColor(0xA0, 0x78, 0xF0),
         ["interviewer_id → Interviewer", "project_id, li_date",
          "pass_fail, supervisor_id → User"]),
        ("CancelledInterview", MUTED,
         ["interviewer_id → Interviewer", "project_id, cancel_date",
          "reason_code, notes"]),
    ]

    cols = 4
    for i, (name, color, fields) in enumerate(tables):
        col, row = i % cols, i // cols
        x = 0.4 + col * 3.22
        y = 1.15 + row * 3.05
        card(s, x, y, 3.0, 2.88, SURFACE)
        add_rect(s, x, y, 3.0, 0.05, color)
        add_text(s, name, x + 0.12, y + 0.1, 2.75, 0.38, size=11, bold=True, color=WHITE)
        for j, f in enumerate(fields):
            add_text(s, "· " + f, x + 0.12, y + 0.58 + j * 0.68, 2.75, 0.62, size=9, color=MUTED)


def slide_dashboard_auth(prs):
    s = blank_slide(prs)
    fill_bg(s)
    slide_header(s, "Dashboard Auth + Deployment", "JWT auth  ·  Fly.io deployment  ·  shared_db.py bridge")

    card(s, 0.5, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 0.5, 1.15, 5.9, 0.06, ACCENT)
    add_text(s, "Authentication Flow", 0.7, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)

    auth_steps = [
        ("Login",        "POST /api/auth/login → returns JWT access token (24h TTL)"),
        ("Token storage", "React: localStorage key ds_auth_token. Sent as Bearer header on every API call."),
        ("Role check",   "FastAPI dependency get_current_user() decodes JWT, checks role against endpoint requirements."),
        ("Dashboard",    "Streamlit sessions use HTTP-only cookie. Same SQLite DB, separate session state."),
        ("shared_db.py", "Bridge module that both the FastAPI backend and Streamlit dashboard import — "
                          "ensures they write to the same SQLite file regardless of working directory."),
    ]
    y = 1.75
    for label, val in auth_steps:
        add_text(s, label, 0.7, y, 2.0, 0.35, size=10, bold=True, color=ACCENT)
        add_text(s, val,   2.75, y, 3.55, 0.75, size=10, color=MUTED)
        y += 0.9

    card(s, 6.8, 1.15, 5.9, 5.85, SURFACE)
    add_rect(s, 6.8, 1.15, 5.9, 0.06, ACCENT2)
    add_text(s, "Deployment", 7.0, 1.28, 5.5, 0.42, size=14, bold=True, color=WHITE)

    deploy = [
        ("QC Engine",      "Render.com — servalab.render.com\nFastAPI backend, React static bundle via Vite build"),
        ("Dashboard",      "Fly.io — dashboard.servallab.com\ndashboard/fly.toml configures app name, region, machine size"),
        ("Database",       "SQLite persisted on Fly.io volume mount.\nNot shared across instances — single-region deployment."),
        ("CI/CD",          "Manual deploy: fly deploy (dashboard) and Render auto-deploy on push to main."),
        ("Env vars",       "VITE_API_URL → FastAPI URL\nSECRET_KEY → JWT signing key\nGROQ_API_KEY → AI verbatim check"),
    ]
    y = 1.75
    for label, val in deploy:
        add_text(s, label, 7.0, y, 2.0, 0.35, size=10, bold=True, color=ACCENT2)
        add_text(s, val,   9.05, y, 3.5, 0.75, size=10, color=MUTED)
        y += 0.9


def slide_closing(prs):
    s = blank_slide(prs)
    fill_bg(s)
    add_rect(s, 0, 0, 13.33, 0.08, ACCENT)
    add_rect(s, 0, 7.42, 13.33, 0.08, ACCENT)
    add_text(s, "Servallab QC Engine", 0.6, 1.9, 12, 0.72, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, "Technical Reference Complete", 0.6, 2.68, 12, 0.52, size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    add_rect(s, 3.5, 3.38, 6.33, 0.02, SURFACE)
    for i, line in enumerate([
        "11 check classes across 5 modules",
        "3 severity levels — Critical / Warning / Info",
        "IQR-based statistical outlier detection",
        "Groq LLM for open-ended verbatim quality scoring",
        "Portable JSON config profiles for project reuse",
        "Manager dashboard — multi-project tracking at dashboard.servallab.com",
    ]):
        add_text(s, "▸  " + line, 2.5, 3.62 + i * 0.52, 8.33, 0.42, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, "github.com/zivodvok2/QC-Engine-desktop", 0.6, 6.6, 12, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ── Build ──────────────────────────────────────────────────────────────────────

def build():
    prs = new_prs()
    slide_title(prs)
    slide_architecture(prs)
    slide_tools(prs)
    slide_pipeline(prs)
    slide_check_overview(prs)
    slide_severity(prs)
    slide_missing(prs)
    slide_range_duration(prs)
    slide_pattern_anomaly(prs)
    slide_logic(prs)
    slide_duplicates(prs)
    slide_straightlining(prs)
    slide_interviewer_duration(prs)
    slide_productivity(prs)
    slide_consent(prs)
    slide_fabrication(prs)
    slide_verbatim(prs)
    slide_verbatim_scoring(prs)
    slide_config_profiles(prs)
    slide_dashboard_overview(prs)
    slide_dashboard_db(prs)
    slide_dashboard_auth(prs)
    slide_closing(prs)
    out = "Servallab_QC_Technical.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")


if __name__ == "__main__":
    build()
